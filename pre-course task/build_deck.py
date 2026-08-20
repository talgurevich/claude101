"""Generates github-page-task-he.pptx — the pre-course assignment deck.

GitHub Pages facts verified against:
https://docs.github.com/en/pages/getting-started-with-github-pages/creating-a-github-pages-site

Re-run after editing:  python3 build_deck.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x1A, 0x18, 0x14)
CREAM = RGBColor(0xF5, 0xF1, 0xEB)
MUTED = RGBColor(0x5C, 0x56, 0x4D)
LIGHT_MUTED = RGBColor(0xC9, 0xC2, 0xB5)
ACCENT = RGBColor(0xD9, 0x77, 0x57)
ACCENT_DARK = RGBColor(0xB8, 0x5F, 0x44)
HAIRLINE = RGBColor(0xE0, 0xDA, 0xD0)
PAPER = RGBColor(0xFA, 0xF7, 0xF2)

BODY = "Arial"
MONO = "Consolas"

MARGIN = 0.5
CONTENT_W = 12.333


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _style_para(para, align, rtl, spc):
    pPr = para._p.get_or_add_pPr()
    pPr.set("algn", {"r": "r", "l": "l", "ctr": "ctr"}[align])
    pPr.set("rtl", "1" if rtl else "0")
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    if spc:
        for run in para.runs:
            run.font._rPr.set("spc", str(spc))


def text(slide, x, y, w, h, runs, size=14, bold=False, color=INK, font=BODY,
         align="r", rtl=True, spc=None, line_spacing=None, anchor=MSO_ANCHOR.TOP):
    """runs: a string, or a list of (text, {overrides}) tuples."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    if isinstance(runs, str):
        runs = [(runs, {})]

    para = tf.paragraphs[0]
    if line_spacing:
        para.line_spacing = line_spacing
    for content, over in runs:
        run = para.add_run()
        run.text = content
        f = run.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)
        if "link" in over:
            run.hyperlink.address = over["link"]
            run.font.color.rgb = over.get("color", color)
        rPr = f._rPr
        for tag in ("a:cs", "a:ea"):
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {})
                rPr.append(el)
            el.set("typeface", over.get("font", font))

    _style_para(para, align, rtl, spc)
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.text_frame.text = ""
    return sp


def eyebrow(slide, label, color=ACCENT):
    text(slide, MARGIN, 0.35, CONTENT_W, 0.3, label, size=11, bold=True, color=color, spc=400)


def heading(slide, title, subtitle=None, size=44):
    text(slide, MARGIN, 0.75, CONTENT_W, 0.85, title, size=size, bold=True, color=INK)
    if subtitle:
        text(slide, MARGIN, 1.66, CONTENT_W, 0.5, subtitle, size=15, color=MUTED, line_spacing=1.25)


def footer_note(slide, label, body, y=6.52):
    rect(slide, MARGIN, y, CONTENT_W, 0.62, fill=PAPER)
    rect(slide, MARGIN + CONTENT_W - 0.04, y, 0.04, 0.62, fill=ACCENT)
    text(slide, MARGIN + 0.25, y + 0.09, CONTENT_W - 0.6, 0.2, label,
         size=9, bold=True, color=ACCENT, spc=300)
    text(slide, MARGIN + 0.25, y + 0.3, CONTENT_W - 0.6, 0.26, body, size=11.5, color=MUTED)


def command_card(slide, x, y, w, h, label, commands, caption=None, dark=True, prompt=False):
    bg = INK if dark else PAPER
    fg = CREAM if dark else INK
    lab = ACCENT if dark else ACCENT_DARK
    rect(slide, x, y, w, h, fill=bg, line=None if dark else HAIRLINE)
    text(slide, x + 0.3, y + 0.24, w - 0.6, 0.22, label, size=10, bold=True, color=lab, spc=300)

    cy = y + 0.6
    for cmd in commands:
        if prompt:
            text(slide, x + 0.3, cy, w - 0.6, 0.32, cmd, size=13, color=fg)
            cy += 0.34
        else:
            text(slide, x + 0.3, cy, w - 0.6, 0.34,
                 [("$ ", {"color": lab}), (cmd, {})],
                 size=13.5, font=MONO, color=fg, align="l", rtl=False)
            cy += 0.42

    if caption:
        text(slide, x + 0.3, y + h - 0.46, w - 0.6, 0.3, caption,
             size=10.5, color=LIGHT_MUTED if dark else MUTED)


def step_rows(slide, y, items, row_h=0.62):
    """Numbered instruction rows, RTL."""
    for i, (num, body) in enumerate(items):
        ry = y + i * row_h
        text(slide, MARGIN + CONTENT_W - 0.42, ry + 0.03, 0.42, 0.28, num,
             size=12, bold=True, color=ACCENT_DARK, font="Georgia", align="l", rtl=False)
        text(slide, MARGIN, ry, CONTENT_W - 0.6, 0.3, body, size=13.5, color=INK)


def slide_new(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if dark:
        set_bg(s, INK)
    return s


# ---------------------------------------------------------------- build

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)


# 01 — cover
s = slide_new(prs, dark=True)
rect(s, -2.5, -1.0, 5.0, 5.0, fill=RGBColor(0x24, 0x21, 0x1C))
rect(s, 11.6, 5.5, 1.2, 0.04, fill=ACCENT)
text(s, MARGIN, 1.55, CONTENT_W, 0.4, "משימת הכנה  ·  לפני המפגש הראשון",
     size=14, bold=True, color=ACCENT, spc=300)
text(s, MARGIN, 2.1, CONTENT_W, 1.2, "עמוד GitHub אישי", size=64, bold=True, color=CREAM)
text(s, MARGIN, 3.6, CONTENT_W, 0.7, "ריפו אחד  ·  עמוד אחד  ·  חי באינטרנט", size=26, color=LIGHT_MUTED)
text(s, MARGIN, 5.7, CONTENT_W, 0.4, "טל גורביץ׳ Tal Gurevich", size=18, bold=True, color=CREAM)
text(s, MARGIN, 6.2, CONTENT_W, 0.3, "המכללה האקדמית תל-חי", size=12, color=MUTED)


# 02 — the deliverable
s = slide_new(prs)
eyebrow(s, "המשימה")
heading(s, "מה צריך להגיש",
        "עמוד אישי אחד, שאתם בונים עם Claude Code, שחי בכתובת ציבורית באינטרנט.")

rect(s, MARGIN, 2.35, CONTENT_W, 1.15, fill=INK)
text(s, MARGIN + 0.3, 2.58, CONTENT_W - 0.6, 0.25, "הכתובת שתגישו",
     size=10, bold=True, color=ACCENT, spc=300)
text(s, MARGIN + 0.3, 2.9, CONTENT_W - 0.6, 0.4, "https://USERNAME.github.io",
     size=24, bold=True, font=MONO, color=CREAM, align="l", rtl=False)

reqs = [
    ("01", "ריפו ציבורי", "בשם המדויק", "USERNAME.github.io", "באותיות קטנות בלבד"),
    ("02", "קובץ index.html", "עמוד אחד", "index.html", "בלי תלויות חיצוניות"),
    ("03", "לפחות שני commits", "מהטרמינל", "git push", "לא מהעורך של GitHub"),
]
cw = (CONTENT_W - 0.6) / 3
for i, (num, title, sub, code, tail) in enumerate(reqs):
    cx = MARGIN + (2 - i) * (cw + 0.3)
    rect(s, cx, 3.72, cw, 2.2, fill=PAPER, line=HAIRLINE)
    rect(s, cx + cw - 0.04, 3.72, 0.04, 2.2, fill=ACCENT)
    text(s, cx + 0.28, 3.98, cw - 0.6, 0.28, num, size=11, bold=True, color=ACCENT_DARK,
         font="Georgia", align="l", rtl=False)
    text(s, cx + 0.28, 3.96, cw - 0.6, 0.34, title, size=17, bold=True, color=INK)
    text(s, cx + 0.28, 4.46, cw - 0.6, 0.28, sub, size=12, color=MUTED)
    text(s, cx + 0.28, 4.86, cw - 0.6, 0.3, code, size=13, bold=True, font=MONO,
         color=ACCENT_DARK, align="l", rtl=False)
    text(s, cx + 0.28, 5.3, cw - 0.6, 0.28, tail, size=11, color=MUTED)

footer_note(s, "המטרה האמיתית",
            "העמוד הוא התירוץ. המטרה שתגיעו למפגש הראשון עם Claude Code מותקן ועם סבב git שלם מאחוריכם.")


# 03 — step 1: setup
s = slide_new(prs)
eyebrow(s, "שלב  ·  01")
heading(s, "הכנה", "לפני שנוגעים ב-GitHub — לוודא ש-Claude Code מותקן ועובד.")

command_card(s, MARGIN, 2.35, CONTENT_W, 1.9, "בדיקה מהירה בטרמינל",
             ["claude --version", "claude doctor"],
             caption="הראשונה מדפיסה מספר גרסה. השנייה נותנת אבחון מלא של ההתקנה.")

step_rows(s, 4.45, [
    ("01", "אם Claude Code עדיין לא מותקן — הדק המלא נמצא בתיקיית installation-instructions בריפו של הקורס."),
    ("02", "צריך חשבון GitHub. אם אין לכם — נרשמים ב-github.com, חינם, לוקח שתי דקות."),
    ("03", "זכרו את שם המשתמש שבחרתם ב-GitHub. הוא יקבע את שם הריפו ואת הכתובת של העמוד."),
])

footer_note(s, "שימו לב",
            "התוכנית החינמית של Claude.ai לא כוללת Claude Code. צריך מנוי Pro / Max, או חשבון Console.")


# 04 — step 2: create the repo
s = slide_new(prs)
eyebrow(s, "שלב  ·  02")
heading(s, "יוצרים את הריפו", "השם של הריפו הוא לא קישוט — הוא מה שגורם ל-GitHub להגיש את העמוד שלכם כאתר.")

step_rows(s, 2.3, [
    ("01", "ב-GitHub: New repository."),
    ("02", "השם חייב להיות בדיוק שם המשתמש שלכם ואחריו github.io — הכול באותיות קטנות."),
    ("03", "מסמנים Public. בחשבון חינמי, ריפו פרטי לא יפורסם כאתר."),
    ("04", "מסמנים Add a README file, ואז Create repository."),
])

command_card(s, MARGIN, 4.85, CONTENT_W, 1.5, "משכפלים למחשב ונכנסים לתיקייה",
             ["git clone https://github.com/USERNAME/USERNAME.github.io.git",
              "cd USERNAME.github.io"])

footer_note(s, "הטעות הכי נפוצה",
            "שם ריפו שגוי או עם אות גדולה — הריפו ייווצר, אבל שום עמוד לא יעלה. תבדקו תו-תו.")


# 05 — step 3: build it with Claude Code
s = slide_new(prs)
eyebrow(s, "שלב  ·  03")
heading(s, "בונים עם Claude Code", "לא כותבים HTML ביד. מריצים claude בתוך התיקייה ונותנים לו לראיין אתכם.")

command_card(s, MARGIN, 2.3, CONTENT_W, 1.0, "מפעילים בתוך התיקייה של הריפו", ["claude"])

command_card(s, MARGIN, 3.45, CONTENT_W, 2.5, "פרומפט לפתיחה  ·  העתיקו והתאימו",
             ["בנה לי עמוד אישי בקובץ index.html יחיד.",
              "לפני שאתה כותב — ראיין אותי: שאל חמש שאלות על מי אני, מה אני לומד,",
              "ומה בניתי. חכה לתשובות שלי, ורק אז כתוב את העמוד.",
              "עברית, RTL, עיצוב נקי, בלי ספריות חיצוניות."],
             caption="שימו לב מה הפרומפט הזה עושה: מגדיר פורמט, מבקש ראיון לפני כתיבה, ומציב אילוצים.",
             prompt=True)

footer_note(s, "למה דווקא ככה",
            "פרומפט שמבקש לשאול לפני שהוא כותב מייצר תוצאה אישית. פרומפט של שורה אחת מייצר תבנית גנרית.")


# 06 — step 4: publish
s = slide_new(prs)
eyebrow(s, "שלב  ·  04")
heading(s, "מעלים לאוויר", "שלוש פקודות, ואז בדיקה אחת בדפדפן.")

command_card(s, MARGIN, 2.3, CONTENT_W, 1.95, "שומרים ודוחפים ל-GitHub",
             ["git add index.html",
              'git commit -m "add personal page"',
              "git push"])

step_rows(s, 4.5, [
    ("01", "ב-GitHub: Settings ← Pages, ומוודאים ש-Source מוגדר לענף main."),
    ("02", "מחכים. פרסום של שינוי יכול לקחת עד עשר דקות — זה נורמלי."),
    ("03", "פותחים את https://USERNAME.github.io בדפדפן. רואים את העמוד? סיימתם."),
])

footer_note(s, "אם העמוד לא עולה",
            "רענון קשיח, או חלון גלישה בסתר. דפדפנים שומרים בזיכרון את שגיאת ה-404 מלפני שהאתר פורסם.")


# 07 — content requirements
s = slide_new(prs)
eyebrow(s, "תוכן")
heading(s, "מה חייב להיות בעמוד", "חמישה דברים. מעבר לזה — תעצבו ותוסיפו מה שבא לכם.")

items = [
    ("מי אתם", "שם, ומשפט אחד שאומר משהו אמיתי עליכם"),
    ("מה אתם לומדים", "שנה, מסלול, ומה מעניין אתכם בו"),
    ("משהו שבניתם", "פרויקט, תרגיל, צד — אפילו קטן. עם משפט הסבר"),
    ("מה אתם רוצים מהקורס", "משפט אחד. זה גם עוזר לי להתאים את המפגשים"),
    ("דרך ליצור קשר", "לינקדאין או מייל שנוח לכם שיהיה פומבי"),
]
y = 2.3
for i, (t, d) in enumerate(items):
    rect(s, MARGIN, y, CONTENT_W, 0.72, fill=PAPER, line=HAIRLINE)
    rect(s, MARGIN + CONTENT_W - 0.04, y, 0.04, 0.72, fill=ACCENT)
    text(s, MARGIN + 0.3, y + 0.11, CONTENT_W - 0.65, 0.3, t, size=14.5, bold=True, color=INK)
    text(s, MARGIN + 0.3, y + 0.4, CONTENT_W - 0.65, 0.28, d, size=11.5, color=MUTED)
    y += 0.82

footer_note(s, "פרטיות — קראו את זה",
            "העמוד פומבי לחלוטין. בלי תעודת זהות, בלי טלפון, בלי כתובת מגורים. מה שעולה לשם, עולה לכולם.",
            y=6.5)


# 08 — submission + resources
s = slide_new(prs, dark=True)
text(s, MARGIN, 0.55, CONTENT_W, 0.3, "הגשה ומשאבים", size=11, bold=True, color=ACCENT, spc=400)
text(s, MARGIN, 0.95, CONTENT_W, 0.8, "איך מגישים", size=40, bold=True, color=CREAM)
text(s, MARGIN, 1.85, CONTENT_W, 0.4,
     "שולחים כתובת אחת. אם היא נפתחת בדפדפן — המשימה הוגשה.",
     size=14, color=LIGHT_MUTED)

rect(s, MARGIN, 2.45, CONTENT_W, 0.95, fill=RGBColor(0x24, 0x21, 0x1C))
text(s, MARGIN + 0.3, 2.68, CONTENT_W - 0.6, 0.45, "https://USERNAME.github.io",
     size=22, bold=True, font=MONO, color=ACCENT, align="l", rtl=False)
text(s, MARGIN + 0.3, 3.08, CONTENT_W - 0.6, 0.25,
     "שולחים במייל, עם שם מלא בגוף ההודעה", size=11.5, color=LIGHT_MUTED)

links = [
    ("התקנת Claude Code — הדק של הקורס", "github.com/talgurevich/claude101"),
    ("התיעוד הרשמי של Claude Code", "code.claude.com/docs/en/setup"),
    ("GitHub Pages — המדריך הרשמי", "docs.github.com/en/pages"),
    ("ללמוד git מאפס", "docs.github.com/en/get-started"),
]
y = 3.75
for label, url in links:
    text(s, MARGIN, y, CONTENT_W * 0.44, 0.3, label, size=12.5, color=CREAM)
    text(s, MARGIN + CONTENT_W * 0.46, y, CONTENT_W * 0.54, 0.3,
         [(url, {"link": "https://" + url, "color": ACCENT})],
         size=12.5, font=MONO, color=ACCENT, align="l", rtl=False)
    rect(s, MARGIN, y + 0.38, CONTENT_W, 0.01, fill=RGBColor(0x33, 0x2F, 0x28))
    y += 0.58

text(s, MARGIN, 6.35, CONTENT_W, 0.4,
     "נתקעתם? אל תאבדו על זה ערב. תכתבו לי, או תביאו את השגיאה למפגש — זה בדיוק מה שנעשה בכיתה.",
     size=12, color=MUTED)


out = "/Users/talgurevich/Documents/Claude101/task/github-page-task-he.pptx"
prs.save(out)
print("wrote", out, "|", len(prs.slides._sldIdLst), "slides")
