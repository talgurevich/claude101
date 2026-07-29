"""Generates claude-code-install-he.pptx — the installation deck.

Source of truth for the commands: https://code.claude.com/docs/en/setup
Re-run after the docs change:  python3 build_deck.py
"""

import copy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x1A, 0x18, 0x14)
CREAM = RGBColor(0xF5, 0xF1, 0xEB)
MUTED = RGBColor(0x5C, 0x56, 0x4D)
LIGHT_MUTED = RGBColor(0xC9, 0xC2, 0xB5)
ACCENT = RGBColor(0xD9, 0x77, 0x57)
ACCENT_DARK = RGBColor(0xB8, 0x5F, 0x44)
HAIRLINE = RGBColor(0xE0, 0xDA, 0xD0)

BODY = "Arial"
MONO = "Consolas"

MARGIN = 0.5
CONTENT_W = 12.333


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _style_para(para, align, rtl, spc):
    pPr = para._p.get_or_add_pPr()
    pPr.set("algn", {"r": "r", "l": "l", "ctr": "ctr"}[align])
    pPr.set("rtl", "1" if rtl else "0")
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    if spc:
        for run in para.runs:
            run.font._rPr.set("spc", str(spc))


def text(
    slide,
    x,
    y,
    w,
    h,
    runs,
    size=14,
    bold=False,
    color=INK,
    font=BODY,
    align="r",
    rtl=True,
    spc=None,
    line_spacing=None,
    anchor=MSO_ANCHOR.TOP,
):
    """runs: a string, or a list of (text, {overrides}) tuples."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    if isinstance(runs, str):
        runs = [(runs, {})]

    para = tf.paragraphs[0]
    if line_spacing:
        para.line_spacing = line_spacing
    for content, over in runs:
        run = para.add_run()
        run.text = content
        f = run.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)
        if "link" in over:
            run.hyperlink.address = over["link"]
            run.font.color.rgb = over.get("color", color)
        # force the complex-script font too, or Hebrew falls back to a default
        rPr = f._rPr
        for tag in ("a:cs", "a:ea"):
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {})
                rPr.append(el)
            el.set("typeface", over.get("font", font))

    _style_para(para, align, rtl, spc)
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.text_frame.text = ""
    return sp


def eyebrow(slide, label, color=ACCENT):
    text(slide, MARGIN, 0.35, CONTENT_W, 0.3, label, size=11, bold=True, color=color, spc=400)


def heading(slide, title, subtitle=None, size=44):
    text(slide, MARGIN, 0.75, CONTENT_W, 0.85, title, size=size, bold=True, color=INK)
    if subtitle:
        text(slide, MARGIN, 1.66, CONTENT_W, 0.5, subtitle, size=15, color=MUTED, line_spacing=1.25)


def footer_note(slide, label, body, y=6.52):
    """Accent-bordered strip at the bottom — the 'why' behind the slide."""
    rect(slide, MARGIN, y, CONTENT_W, 0.62, fill=RGBColor(0xFA, 0xF7, 0xF2))
    rect(slide, MARGIN + CONTENT_W - 0.04, y, 0.04, 0.62, fill=ACCENT)
    text(slide, MARGIN + 0.25, y + 0.09, CONTENT_W - 0.6, 0.2, label,
         size=9, bold=True, color=ACCENT, spc=300)
    text(slide, MARGIN + 0.25, y + 0.3, CONTENT_W - 0.6, 0.26, body, size=11.5, color=MUTED)


def command_card(slide, x, y, w, h, label, commands, caption=None, dark=True):
    """A terminal block. Commands render LTR in mono, everything else RTL."""
    bg = INK if dark else RGBColor(0xF5, 0xF1, 0xEB)
    fg = CREAM if dark else INK
    lab = ACCENT if dark else ACCENT_DARK
    rect(slide, x, y, w, h, fill=bg, line=None if dark else HAIRLINE)

    text(slide, x + 0.3, y + 0.26, w - 0.6, 0.22, label, size=10, bold=True, color=lab, spc=300)

    cy = y + 0.62
    for cmd in commands:
        text(slide, x + 0.3, cy, w - 0.6, 0.34, [("$ ", {"color": lab}), (cmd, {})],
             size=13.5, font=MONO, color=fg, align="l", rtl=False)
        cy += 0.42

    if caption:
        text(slide, x + 0.3, y + h - 0.48, w - 0.6, 0.3, caption,
             size=10.5, color=LIGHT_MUTED if dark else MUTED, line_spacing=1.2)


def slide_new(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if dark:
        set_bg(s, INK)
    return s


# ---------------------------------------------------------------- build

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)


# 01 — cover
s = slide_new(prs, dark=True)
rect(s, -2.5, -1.0, 5.0, 5.0, fill=RGBColor(0x24, 0x21, 0x1C))
rect(s, 11.6, 5.5, 1.2, 0.04, fill=ACCENT)
text(s, MARGIN, 1.6, CONTENT_W, 0.4, "מדריך התקנה", size=14, bold=True, color=ACCENT, spc=300)
text(s, MARGIN, 2.15, CONTENT_W, 1.4, "Claude Code", size=72, bold=True, color=CREAM, align="l", rtl=False)
text(s, MARGIN, 3.7, CONTENT_W, 0.7, "התקנה בטרמינל  ·  macOS · Windows · Linux",
     size=26, color=LIGHT_MUTED)
text(s, MARGIN, 5.7, CONTENT_W, 0.4, "טל גורביץ׳ Tal Gurevich", size=18, bold=True, color=CREAM)
text(s, MARGIN, 6.2, CONTENT_W, 0.3, "המכללה האקדמית תל-חי", size=12, color=MUTED)


# 02 — prerequisites
s = slide_new(prs)
eyebrow(s, "לפני שמתחילים")
heading(s, "דרישות מערכת", "בדקו את ארבעת אלה לפני שאתם מריצים פקודה כלשהי. רוב תקלות ההתקנה מתחילות כאן.")

cards = [
    ("01", "מערכת הפעלה", "macOS 13+  ·  Windows 10 1809+  ·  Ubuntu 20.04+", "Debian 10+  ·  Alpine Linux 3.19+"),
    ("02", "חומרה ורשת", "‏4GB RAM ומעלה  ·  מעבד x64 או ARM64", "חיבור אינטרנט פעיל, ממדינה נתמכת"),
    ("03", "חשבון", "‏Pro · Max · Team · Enterprise · Console", "התוכנית החינמית של Claude.ai לא כוללת גישה"),
    ("04", "טרמינל", "‏Bash · Zsh · PowerShell · CMD", "‏ripgrep מגיע מובנה עם Claude Code"),
]
cw, ch = (CONTENT_W - 0.3) / 2, 1.7
for i, (num, title, l1, l2) in enumerate(cards):
    # RTL reading order: card 01 sits top-right
    cx = MARGIN + (1 - i % 2) * (cw + 0.3)
    cy = 2.35 + (i // 2) * (ch + 0.3)
    rect(s, cx, cy, cw, ch, fill=RGBColor(0xFA, 0xF7, 0xF2), line=HAIRLINE)
    text(s, cx + 0.3, cy + 0.28, cw - 0.6, 0.28, num, size=11, bold=True, color=ACCENT_DARK,
         font="Georgia", align="l", rtl=False)
    text(s, cx + 0.3, cy + 0.26, cw - 0.6, 0.34, title, size=17, bold=True, color=INK)
    text(s, cx + 0.3, cy + 0.78, cw - 0.6, 0.3, l1, size=12, color=MUTED)
    text(s, cx + 0.3, cy + 1.15, cw - 0.6, 0.3, l2, size=12, color=MUTED)

footer_note(s, "למה זה חשוב",
            "התוכנית החינמית של Claude.ai לא כוללת Claude Code. אם ההתחברות נכשלת — זו הסיבה הראשונה לבדוק.")


# 03 — macOS
s = slide_new(prs)
eyebrow(s, "מערכת הפעלה  ·  01")
heading(s, "macOS", "פותחים את Terminal (או iTerm) ומדביקים שורה אחת. אין צורך ב-sudo.")

command_card(s, MARGIN, 2.3, CONTENT_W, 1.75,
             "ההתקנה המומלצת  ·  NATIVE INSTALLER",
             ["curl -fsSL https://claude.ai/install.sh | bash"],
             caption="מתעדכן אוטומטית ברקע. זו הדרך שאנחנו ממליצים עליה בקורס.")

hw = (CONTENT_W - 0.3) / 2
command_card(s, MARGIN + hw + 0.3, 4.25, hw, 2.0, "חלופה  ·  HOMEBREW",
             ["brew install --cask claude-code"],
             caption="לא מתעדכן לבד. שדרוג ידני:  brew upgrade claude-code", dark=False)
command_card(s, MARGIN, 4.25, hw, 2.0, "הפעלה ראשונה",
             ["claude"],
             caption="נפתח דפדפן להתחברות, ואז נפתח סשן בתיקייה הנוכחית.", dark=False)

footer_note(s, "טיפ",
            "‏claude רץ בתיקייה שאתם נמצאים בה. תמיד cd לתוך הפרויקט לפני שמפעילים — זה מה שהוא יראה.")


# 04 — Windows
s = slide_new(prs)
eyebrow(s, "מערכת הפעלה  ·  02")
heading(s, "Windows", "שתי פקודות שונות — תלוי אם אתם ב-PowerShell או ב-CMD. אם השורה מתחילה ב-‎PS C:\\‎ אתם ב-PowerShell.")

command_card(s, MARGIN, 2.35, CONTENT_W, 1.35, "POWERSHELL",
             ["irm https://claude.ai/install.ps1 | iex"])
command_card(s, MARGIN, 3.85, CONTENT_W, 1.35, "CMD",
             ["curl -fsSL https://claude.ai/install.cmd -o install.cmd && install.cmd && del install.cmd"])

hw = (CONTENT_W - 0.3) / 2
command_card(s, MARGIN + hw + 0.3, 5.35, hw, 1.0, "חלופה  ·  WINGET",
             ["winget install Anthropic.ClaudeCode"], dark=False)
command_card(s, MARGIN, 5.35, hw, 1.0, "‏WSL — מריצים את פקודת ה-LINUX",
             ["curl -fsSL https://claude.ai/install.sh | bash"], dark=False)

footer_note(s, "שווה לדעת",
            "התקינו Git for Windows כדי ש-Claude Code ישתמש ב-Bash. בלעדיו הוא נופל ל-PowerShell. סנדבוקסינג עובד רק ב-WSL 2.",
            y=6.55)


# 05 — Linux
s = slide_new(prs)
eyebrow(s, "מערכת הפעלה  ·  03")
heading(s, "Linux", "אותה פקודה כמו ב-macOS. למי שרוצה שהעדכונים יגיעו דרך מנהל החבילות — יש ריפוזיטוריז חתומים.")

command_card(s, MARGIN, 2.3, CONTENT_W, 1.5, "ההתקנה המומלצת  ·  NATIVE INSTALLER",
             ["curl -fsSL https://claude.ai/install.sh | bash"])

tw = (CONTENT_W - 0.6) / 3
specs = [
    ("DEBIAN / UBUNTU  ·  APT", ["sudo apt install claude-code"], "אחרי רישום המפתח והריפו"),
    ("FEDORA / RHEL  ·  DNF", ["sudo dnf install claude-code"], "אחרי יצירת claude-code.repo"),
    ("ALPINE  ·  APK", ["apk add claude-code"], "דורש bash curl libgcc libstdc++ ripgrep"),
]
for i, (lab, cmds, cap) in enumerate(specs):
    command_card(s, MARGIN + (2 - i) * (tw + 0.3), 4.0, tw, 1.75, lab, cmds, caption=cap, dark=False)

footer_note(s, "שימו לב",
            "התקנה דרך מנהל חבילות לא מתעדכנת אוטומטית — העדכון מגיע עם עדכוני המערכת הרגילים שלכם.",
            y=6.0)
text(s, MARGIN, 6.75, CONTENT_W, 0.3,
     "מפתח החתימה: 31DD DE24 DDFA B679 F42D 7BD2 BAA9 29FF 1A7E CACE",
     size=9.5, font=MONO, color=MUTED, align="l", rtl=False)


# 06 — verify + authenticate
s = slide_new(prs)
eyebrow(s, "אחרי ההתקנה")
heading(s, "אימות והתחברות", "שלוש פקודות שאומרות לכם אם ההתקנה באמת הצליחה — לפני שאתם מתחילים לעבוד.")

rows = [
    ("claude --version", "מדפיס מספר גרסה, למשל 2.1.211 (Claude Code). אם אתם מקבלים command not found — ההתקנה לא נכנסה ל-PATH."),
    ("claude doctor", "אבחון קריאה-בלבד: תקינות ההתקנה, שגיאות בקובצי ההגדרות, ואזהרות עם הצעות תיקון. לא פותח סשן."),
    ("claude", "מפעיל סשן אינטראקטיבי. בהרצה הראשונה נפתח דפדפן להתחברות לחשבון Anthropic שלכם."),
]
y = 2.35
for cmd, desc in rows:
    rect(s, MARGIN, y, CONTENT_W, 1.15, fill=RGBColor(0xFA, 0xF7, 0xF2), line=HAIRLINE)
    rect(s, MARGIN + CONTENT_W - 0.04, y, 0.04, 1.15, fill=ACCENT)
    text(s, MARGIN + 0.32, y + 0.22, CONTENT_W - 0.7, 0.32, [("$ ", {"color": ACCENT_DARK}), (cmd, {})],
         size=15, bold=True, font=MONO, color=INK, align="l", rtl=False)
    text(s, MARGIN + 0.32, y + 0.66, CONTENT_W - 0.7, 0.32, desc, size=12, color=MUTED)
    y += 1.32

footer_note(s, "אם claude doctor מתלונן",
            "אל תתחילו לנחש. קחו את ההודעה המדויקת לדף troubleshoot-install — הוא בנוי כטבלת שגיאה ← תיקון.")


# 07 — login options
s = slide_new(prs)
eyebrow(s, "התחברות ראשונה")
heading(s, "איך מתחברים",
        "בהרצה הראשונה של claude נפתח דפדפן. לא נפתח? לחצו c להעתקת הקישור. הדפדפן מציג קוד (נפוץ ב-WSL2 ו-SSH)? הדביקו אותו בטרמינל.")

options = [
    ("01", "מנוי Claude.ai", "Pro  ·  Max  ·  Team  ·  Enterprise",
     "מתחברים עם חשבון claude.ai הרגיל. זו האפשרות של רוב",
     "הסטודנטים: משלמים מנוי חודשי קבוע, בלי חיוב לפי שימוש."),
    ("02", "חשבון Claude Console", "חיוב לפי שימוש  ·  API",
     "מתחברים עם פרטי ה-Console. מנהל הארגון צריך להזמין",
     "אתכם קודם. החיוב הוא לפי טוקנים שנצרכו בפועל."),
]
ow = (CONTENT_W - 0.3) / 2
for i, (num, title, tag, d1, d2) in enumerate(options):
    ox = MARGIN + (1 - i) * (ow + 0.3)
    rect(s, ox, 2.42, ow, 1.95, fill=RGBColor(0xFA, 0xF7, 0xF2), line=HAIRLINE)
    rect(s, ox + ow - 0.04, 2.42, 0.04, 1.95, fill=ACCENT)
    text(s, ox + 0.3, 2.68, ow - 0.65, 0.28, num, size=11, bold=True, color=ACCENT_DARK,
         font="Georgia", align="l", rtl=False)
    text(s, ox + 0.3, 2.66, ow - 0.65, 0.34, title, size=19, bold=True, color=INK)
    text(s, ox + 0.3, 3.18, ow - 0.65, 0.3, tag, size=11.5, bold=True, color=ACCENT_DARK, font=MONO)
    text(s, ox + 0.3, 3.6, ow - 0.65, 0.3, d1, size=12, color=MUTED)
    text(s, ox + 0.3, 3.9, ow - 0.65, 0.3, d2, size=12, color=MUTED)

rect(s, MARGIN, 4.55, CONTENT_W, 0.85, fill=RGBColor(0xFA, 0xF7, 0xF2), line=HAIRLINE)
rect(s, MARGIN + CONTENT_W - 0.04, 4.55, 0.04, 0.85, fill=LIGHT_MUTED)
text(s, MARGIN + 0.3, 4.74, CONTENT_W - 0.65, 0.3,
     [("03  ", {"font": "Georgia", "bold": True, "color": ACCENT_DARK, "size": 11}),
      ("ספקי ענן — Amazon Bedrock · Google Vertex · Microsoft Foundry", {})],
     size=14, bold=True, color=INK)
text(s, MARGIN + 0.3, 5.06, CONTENT_W - 0.65, 0.3,
     "בוחרים 3rd-party platform במסך ההתחברות, או מגדירים משתני סביבה מראש. אין התחברות דרך דפדפן.",
     size=11.5, color=MUTED)

cmds = [
    ("/login", "התחברות או החלפת חשבון"),
    ("/logout", "התנתקות ואיפוס ההגדרה"),
    ("/status", "מי מחובר ובאיזו שיטה"),
]
rect(s, MARGIN, 5.58, CONTENT_W, 0.78, fill=INK)
iw = (CONTENT_W - 0.6) / 3
for i, (cmd, gloss) in enumerate(cmds):
    ix = MARGIN + 0.3 + (2 - i) * iw
    text(s, ix, 5.75, iw - 0.2, 0.28, cmd, size=14, bold=True, color=ACCENT, font=MONO,
         align="l", rtl=False)
    text(s, ix, 6.04, iw - 0.2, 0.26, gloss, size=10.5, color=LIGHT_MUTED, align="l")

footer_note(s, "המלכודת הנפוצה",
            "אם ANTHROPIC_API_KEY מוגדר בסביבה — הוא גובר על המנוי שלכם. לחזרה למנוי: unset ANTHROPIC_API_KEY, ואז /status לאימות.",
            y=6.5)


# 08 — troubleshooting
s = slide_new(prs)
eyebrow(s, "כשזה לא עובד")
heading(s, "תקלות נפוצות", "ארבע השגיאות שרוב הכיתה נתקלת בהן. אל תאבדו שעה על התקנה — זה לא הנושא של הקורס.")

issues = [
    ("command not found: claude",
     "ה-PATH לא עודכן בטרמינל הפתוח. פתחו חלון טרמינל חדש, או הריצו source ~/.zshrc (Mac) / source ~/.bashrc (Linux)."),
    ("‏403 או syntax error near unexpected token '<'",
     "ה-curl קיבל דף שגיאה במקום סקריפט. בדקו רשת/פרוקסי, ואז עברו לדף troubleshoot-install לשיטת התקנה חלופית."),
    ("‏The token '&&' is not a valid statement separator",
     "הרצתם את פקודת ה-CMD בתוך PowerShell. אלה שתי פקודות שונות — בדקו אם השורה מתחילה ב-‎PS C:\\‎."),
    ("עברית מוצגת כריבועים או כסימני שאלה",
     "פונט הטרמינל לא תומך בעברית. ב-Mac: Terminal ← Settings ← Profiles ← Text, ובחרו Menlo או SF Mono."),
]
y = 2.3
for title, fix in issues:
    rect(s, MARGIN, y, CONTENT_W, 1.0, fill=RGBColor(0xFA, 0xF7, 0xF2), line=HAIRLINE)
    text(s, MARGIN + 0.32, y + 0.18, CONTENT_W - 0.7, 0.3, title, size=13.5, bold=True, color=ACCENT_DARK)
    text(s, MARGIN + 0.32, y + 0.56, CONTENT_W - 0.7, 0.32, fix, size=11.5, color=MUTED)
    y += 1.1

footer_note(s, "כלל אצבע",
            "אף פעם לא sudo npm install -g. זה מייצר בעיות הרשאות ומסכן אבטחה — וזו סיבה שכיחה להתקנה שבורה.")


# 09 — resources
s = slide_new(prs, dark=True)
text(s, MARGIN, 0.55, CONTENT_W, 0.3, "משאבים", size=11, bold=True, color=ACCENT, spc=400)
text(s, MARGIN, 0.95, CONTENT_W, 0.8, "כל הקישורים הרשמיים", size=40, bold=True, color=CREAM)
text(s, MARGIN, 1.85, CONTENT_W, 0.4,
     "התיעוד הרשמי של Anthropic הוא מקור האמת. הוא משתנה — תמיד עדיף עליו מאשר על סלייד.",
     size=14, color=LIGHT_MUTED)

links = [
    ("התקנה מלאה — דרישות, כל שיטות ההתקנה, הסרה", "code.claude.com/docs/en/setup"),
    ("התחלה מהירה — הסשן הראשון שלכם", "code.claude.com/docs/en/quickstart"),
    ("פתרון בעיות התקנה והתחברות", "code.claude.com/docs/en/troubleshoot-install"),
    ("מדריך טרמינל — למי שלא עבד בטרמינל קודם", "code.claude.com/docs/en/terminal-guide"),
    ("אימות וחשבונות — כולל Bedrock ו-Vertex", "code.claude.com/docs/en/authentication"),
    ("אינדקס כל התיעוד (קובץ אחד, נוח לסוכן)", "code.claude.com/docs/llms.txt"),
]
y = 2.5
for label, url in links:
    text(s, MARGIN, y, CONTENT_W * 0.44, 0.3, label, size=12.5, color=CREAM)
    text(s, MARGIN + CONTENT_W * 0.46, y, CONTENT_W * 0.54, 0.3,
         [(url, {"link": "https://" + url, "color": ACCENT})],
         size=12.5, font=MONO, color=ACCENT, align="l", rtl=False)
    rect(s, MARGIN, y + 0.38, CONTENT_W, 0.01, fill=RGBColor(0x33, 0x2F, 0x28))
    y += 0.58

text(s, MARGIN, 6.35, CONTENT_W, 0.5,
     "קישורים ישנים מסוג docs.claude.com/en/docs/claude-code/... עדיין עובדים, אבל מבצעים הפניה (301) אל code.claude.com/docs/en/...",
     size=10.5, color=MUTED, line_spacing=1.2)


out = "/Users/talgurevich/Documents/Claude101/installation-instructions/claude-code-install-he.pptx"
prs.save(out)
print("wrote", out, "|", len(prs.slides._sldIdLst), "slides")
