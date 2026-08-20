"""Generates final-project-deck-he.pptx — the capstone assignment deck.

Content is the presentable form of final-project/README.md. If the README
changes, change this too — the README stays the source of truth.

Re-run after editing:  python3 build_deck.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x1A, 0x18, 0x14)
INK_SOFT = RGBColor(0x24, 0x21, 0x1C)
CREAM = RGBColor(0xF5, 0xF1, 0xEB)
MUTED = RGBColor(0x5C, 0x56, 0x4D)
LIGHT_MUTED = RGBColor(0xC9, 0xC2, 0xB5)
DIM = RGBColor(0x9B, 0x95, 0x88)
ACCENT = RGBColor(0xD9, 0x77, 0x57)
ACCENT_DARK = RGBColor(0xB8, 0x5F, 0x44)
HAIRLINE = RGBColor(0xE0, 0xDA, 0xD0)
PAPER = RGBColor(0xFA, 0xF7, 0xF2)
HAIRLINE_DARK = RGBColor(0x33, 0x2F, 0x28)

BODY = "Arial"
MONO = "Consolas"

MARGIN = 0.5
CONTENT_W = 12.333


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _style_para(para, align, rtl, spc):
    pPr = para._p.get_or_add_pPr()
    pPr.set("algn", {"r": "r", "l": "l", "ctr": "ctr"}[align])
    pPr.set("rtl", "1" if rtl else "0")
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    if spc:
        for run in para.runs:
            run.font._rPr.set("spc", str(spc))


def text(slide, x, y, w, h, runs, size=14, bold=False, color=INK, font=BODY,
         align="r", rtl=True, spc=None, line_spacing=None, anchor=MSO_ANCHOR.TOP):
    """runs: a string, or a list of (text, {overrides}) tuples."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    if isinstance(runs, str):
        runs = [(runs, {})]

    para = tf.paragraphs[0]
    if line_spacing:
        para.line_spacing = line_spacing
    for content, over in runs:
        run = para.add_run()
        run.text = content
        f = run.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)
        if "link" in over:
            run.hyperlink.address = over["link"]
            run.font.color.rgb = over.get("color", color)
        rPr = f._rPr
        for tag in ("a:cs", "a:ea"):
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {})
                rPr.append(el)
            el.set("typeface", over.get("font", font))

    _style_para(para, align, rtl, spc)
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.text_frame.text = ""
    return sp


def eyebrow(slide, label, color=ACCENT):
    text(slide, MARGIN, 0.35, CONTENT_W, 0.3, label, size=11, bold=True, color=color, spc=400)


def heading(slide, title, subtitle=None, size=40, color=INK, sub_color=MUTED):
    text(slide, MARGIN, 0.72, CONTENT_W, 0.85, title, size=size, bold=True, color=color)
    if subtitle:
        text(slide, MARGIN, 1.62, CONTENT_W, 0.5, subtitle, size=14.5, color=sub_color,
             line_spacing=1.25)


def footer_note(slide, label, body, y=6.52):
    rect(slide, MARGIN, y, CONTENT_W, 0.62, fill=PAPER)
    rect(slide, MARGIN + CONTENT_W - 0.04, y, 0.04, 0.62, fill=ACCENT)
    text(slide, MARGIN + 0.25, y + 0.09, CONTENT_W - 0.6, 0.2, label,
         size=9, bold=True, color=ACCENT, spc=300)
    text(slide, MARGIN + 0.25, y + 0.3, CONTENT_W - 0.6, 0.26, body, size=11.5, color=MUTED)


def code_card(slide, x, y, w, h, label, lines, caption=None, dark=True, size=13.5):
    bg = INK if dark else PAPER
    fg = CREAM if dark else INK
    lab = ACCENT if dark else ACCENT_DARK
    rect(slide, x, y, w, h, fill=bg, line=None if dark else HAIRLINE)
    text(slide, x + 0.3, y + 0.22, w - 0.6, 0.22, label, size=10, bold=True, color=lab, spc=300)
    cy = y + 0.58
    for ln in lines:
        text(slide, x + 0.3, cy, w - 0.6, 0.32, ln, size=size, font=MONO, color=fg,
             align="l", rtl=False)
        cy += 0.36
    if caption:
        text(slide, x + 0.3, y + h - 0.45, w - 0.6, 0.3, caption,
             size=10.5, color=LIGHT_MUTED if dark else MUTED)


def bullet_card(slide, x, y, w, h, label, title, lines, dark=False, tail=None):
    bg = INK if dark else PAPER
    fg = CREAM if dark else INK
    lab = ACCENT if dark else ACCENT_DARK
    sub = LIGHT_MUTED if dark else MUTED
    rect(slide, x, y, w, h, fill=bg, line=None if dark else HAIRLINE)
    rect(slide, x + w - 0.04, y, 0.04, h, fill=ACCENT if dark else HAIRLINE)
    text(slide, x + 0.28, y + 0.26, w - 0.56, 0.22, label, size=10, bold=True, color=lab, spc=300)
    text(slide, x + 0.28, y + 0.58, w - 0.56, 0.34, title, size=18, bold=True, color=fg)
    cy = y + 1.12
    for ln in lines:
        text(slide, x + 0.28 + 0.22, cy, w - 0.78, 0.3, ln, size=12.5, color=sub,
             line_spacing=1.15)
        text(slide, x + w - 0.42, cy + 0.03, 0.16, 0.2, "·", size=15, bold=True, color=ACCENT)
        cy += 0.42
    if tail:
        text(slide, x + 0.28, y + h - 0.5, w - 0.56, 0.32, tail, size=11.5, bold=True, color=lab)


def row_list(slide, y, items, row_h=0.82, box_h=0.72):
    """Full-width rows: bold title on the right, muted description under it."""
    for i, (t, d) in enumerate(items):
        ry = y + i * row_h
        rect(slide, MARGIN, ry, CONTENT_W, box_h, fill=PAPER, line=HAIRLINE)
        rect(slide, MARGIN + CONTENT_W - 0.04, ry, 0.04, box_h, fill=ACCENT)
        text(slide, MARGIN + 0.3, ry + 0.11, CONTENT_W - 0.65, 0.3, t, size=14.5, bold=True, color=INK)
        text(slide, MARGIN + 0.3, ry + 0.4, CONTENT_W - 0.65, 0.28, d, size=11.5, color=MUTED)


def table_rtl(slide, y, headers, rows, widths, row_h=0.52, mono_cols=(), sizes=None,
              head_h=0.42):
    """RTL table. headers[0] / widths[0] is the RIGHTMOST column."""
    sizes = sizes or [12.5] * len(headers)
    xs, x = [], MARGIN + CONTENT_W
    for w in widths:
        x -= w
        xs.append(x)

    for i, head in enumerate(headers):
        mono = i in mono_cols
        text(slide, xs[i] + 0.16, y + 0.08, widths[i] - 0.32, 0.26, head,
             size=10, bold=True, color=ACCENT_DARK, spc=200,
             align="l" if mono else "r", rtl=not mono)
    rect(slide, MARGIN, y + head_h - 0.02, CONTENT_W, 0.015, fill=ACCENT)

    ry = y + head_h + 0.06
    for r, row in enumerate(rows):
        if r % 2 == 0:
            rect(slide, MARGIN, ry, CONTENT_W, row_h, fill=PAPER)
        for i, cell in enumerate(row):
            mono = i in mono_cols
            text(slide, xs[i] + 0.16, ry + (row_h - 0.28) / 2, widths[i] - 0.32, 0.3, cell,
                 size=sizes[i], color=INK if i == 0 else MUTED,
                 bold=(i == 0), font=MONO if mono else BODY,
                 align="l" if mono else "r", rtl=not mono)
        ry += row_h
    return ry


def slide_new(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if dark:
        set_bg(s, INK)
    return s


# ---------------------------------------------------------------- build

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)


# 01 — cover
s = slide_new(prs, dark=True)
rect(s, -2.6, -1.2, 5.2, 5.2, fill=INK_SOFT)
rect(s, MARGIN + CONTENT_W - 1.2, 4.62, 1.2, 0.04, fill=ACCENT)

text(s, MARGIN, 1.35, CONTENT_W, 0.4, "פרויקט גמר  ·  עבודת צוות  ·  2–3 ימי עבודה",
     size=14, bold=True, color=ACCENT, spc=300)
text(s, MARGIN, 1.92, CONTENT_W, 1.2, "צ׳אט מעל מסמכי החוג", size=58, bold=True, color=CREAM)
text(s, MARGIN, 3.45, CONTENT_W, 0.7, "מערכת שעונה רק ממה שכתוב — ומודה כשאין לה תשובה",
     size=24, color=LIGHT_MUTED)

text(s, MARGIN, 5.0, CONTENT_W, 0.35, "Agentic Retrieval  ·  Claude SDK  ·  SQLite FTS5  ·  Streamlit",
     size=15, color=CREAM)
text(s, MARGIN, 5.45, CONTENT_W, 0.3, "הגשה: ריפו ציבורי ב-GitHub", size=12.5, color=DIM)

text(s, MARGIN, 6.3, CONTENT_W, 0.4, "טל גורביץ׳", size=18, bold=True, color=CREAM)
text(s, MARGIN, 6.74, CONTENT_W, 0.3, "המכללה האקדמית תל-חי", size=11.5, color=DIM)


# 02 — what we are building
s = slide_new(prs)
eyebrow(s, "הפרויקט")
heading(s, "מה בונים",
        "ממשק צ׳אט בעברית שעונה על שאלות אמיתיות של סטודנטים מתוך שני מסמכים רשמיים של המכללה — "
        "השנתון של החוג למדעי המחשב, ותקנון הלימודים לתואר ראשון.")

rect(s, MARGIN, 2.55, CONTENT_W, 1.95, fill=INK)
text(s, MARGIN + 0.35, 2.8, CONTENT_W - 0.7, 0.25, "החוזה של המערכת",
     size=10, bold=True, color=ACCENT, spc=300)
contract = [
    "עונה רק ממה שכתוב במסמכים — לא מידע כללי, לא ניחוש",
    "מציינת בכל תשובה מאיפה היא לקחה אותה",
    "אומרת בפירוש ״לא מופיע במסמכים״ כשהתשובה לא שם",
]
cy = 3.2
for line in contract:
    text(s, MARGIN + 0.7, cy, CONTENT_W - 1.1, 0.32, line, size=15, color=CREAM)
    text(s, MARGIN + CONTENT_W - 0.62, cy + 0.02, 0.2, 0.24, "·", size=18, bold=True, color=ACCENT)
    cy += 0.42

text(s, MARGIN, 4.85, CONTENT_W, 0.4, "זה לא chatbot כללי.", size=17, bold=True, color=INK)
text(s, MARGIN, 5.32, CONTENT_W, 0.7,
     "השורה השלישית היא החלק הכי חשוב בפרויקט. מערכת שממציאה תשובה על תקנון לימודים "
     "גרועה בהרבה ממערכת שפשוט לא עונה.",
     size=14.5, color=MUTED, line_spacing=1.25)

footer_note(s, "למה זה מעניין הנדסית",
            "הקורפוס קטן, המסמכים מובנים, והשפה עברית. שלושת אלה משנים לגמרי איזו ארכיטקטורה נכונה כאן.")


# 03 — the corpus
s = slide_new(prs)
eyebrow(s, "חומרי הגלם")
heading(s, "מה אתם מקבלים", "שני PDF-ים אמיתיים בתיקייה final-project/pdfs/ — טקסט אמיתי, לא סריקות.")

docs = [
    ("שנתון תשפ״ז — מדעי המחשב", "45 עמודים", "טבלאות קורסים, נ״ז, דרישות קדם, מסלולים"),
    ("תקנון לתואר ראשון — תשפ״ו", "13 עמודים", "סעיפים ממוספרים: משך לימודים, פריסה, תנאים"),
]
cw = (CONTENT_W - 0.35) / 2
for i, (title, pages, desc) in enumerate(docs):
    cx = MARGIN + (1 - i) * (cw + 0.35)
    rect(s, cx, 2.5, cw, 1.55, fill=PAPER, line=HAIRLINE)
    rect(s, cx + cw - 0.04, 2.5, 0.04, 1.55, fill=ACCENT)
    text(s, cx + 0.3, 2.74, cw - 0.6, 0.34, title, size=17, bold=True, color=INK)
    text(s, cx + 0.3, 3.18, cw - 0.6, 0.3, pages, size=13, bold=True, color=ACCENT_DARK)
    text(s, cx + 0.3, 3.55, cw - 0.6, 0.3, desc, size=11.5, color=MUTED)

stats = [("58", "עמודים בסך הכל"), ("233K", "תווים של טקסט"), ("0", "צורך ב-OCR")]
sw = (CONTENT_W - 0.7) / 3
for i, (num, lab) in enumerate(stats):
    sx = MARGIN + (2 - i) * (sw + 0.35)
    rect(s, sx, 4.3, sw, 1.1, fill=INK)
    text(s, sx + 0.3, 4.5, sw - 0.6, 0.45, num, size=28, bold=True, color=ACCENT,
         font=MONO, align="l", rtl=False)
    text(s, sx + 0.3, 4.98, sw - 0.6, 0.28, lab, size=12, color=LIGHT_MUTED)

text(s, MARGIN, 5.65, CONTENT_W, 0.5,
     "חילוץ נאיבי של הטקסט לא יעבוד. תגלו את זה בעצמכם בשלב 1 — אל תבנו שכבה שלמה מעל פלט שלא הסתכלתם עליו.",
     size=14, bold=True, color=INK)

footer_note(s, "שער מעבר",
            "לפני שממשיכים משלב החילוץ: פותחים את הטקסט שיצא, וקוראים אותו בעיניים. לא מדגם — באמת קוראים.")


# 04 — the stack
s = slide_new(prs)
eyebrow(s, "החלטות שכבר קיבלנו בשבילכם")
heading(s, "הסטאק המומלץ",
        "המסלול הזה נבדק והוא עובד. אתם לא חייבים אותו — אבל אם בחרתם אחרת, תכתבו ב-README משפט שמסביר למה.")

table_rtl(s, 2.45,
          ["שכבה", "ברירת מחדל", "הערה"],
          [
              ["שפה", "Python 3.11+", "כל השאר עומד על זה"],
              ["חילוץ PDF", "pdfplumber", "pip בלבד, עובד זהה ב-Mac וב-Windows"],
              ["טיפול ב-RTL", "python-bidi", "להפוך רק את מה שצריך — לא את כל המחרוזת"],
              ["אינדקס חיפוש", "SQLite FTS5", "דירוג bm25() מובנה, טוקנייזר unicode61"],
              ["מודל", "anthropic · claude-sonnet-5", "לולאת tool use"],
              ["ממשק", "Streamlit", "הרכיבים st.chat_message ו-st.chat_input"],
              ["הרצה", "מקומית בלבד", "לא מעלים לאוויר — ראו שקופית המפתחות"],
          ],
          widths=[2.1, 3.6, 6.633],
          mono_cols=(1,),
          sizes=[13, 12.5, 12])

footer_note(s, "אם אתם מתלבטים",
            "קחו את ברירת המחדל ותמשיכו הלאה. בחירת ספריות היא לא החלק המעניין בפרויקט הזה.")


# 05 — what to avoid
s = slide_new(prs)
eyebrow(s, "החלטה נגדית")
heading(s, "מה שכדאי מאוד להימנע ממנו",
        "לא בגלל שאלה כלים גרועים — בגלל שבקורפוס בגודל הזה הם תשתית מיותרת.")

rect(s, MARGIN, 2.5, CONTENT_W, 0.95, fill=INK)
text(s, MARGIN + 0.35, 2.72, CONTENT_W - 0.7, 0.22, "לא בפרויקט הזה",
     size=10, bold=True, color=ACCENT, spc=300)
text(s, MARGIN + 0.35, 3.0, CONTENT_W - 0.7, 0.34,
     "embeddings   ·   vector database   ·   Postgres   ·   pgvector   ·   Chroma   ·   torch",
     size=15, font=MONO, color=CREAM, align="l", rtl=False)

row_list(s, 3.75, [
    ("אחזור סמנטי בעברית הוא קרב שיאכל לכם יום שלם מתוך שלושה",
     "והרווח עליו, ב-58 עמודים של טקסט מובנה, קרוב לאפס."),
    ("האינדקס שלכם הוא קובץ אחד בגודל כמה מאות קילובייט",
     "זה בסדר גמור. קובץ SQLite אחד, בלי שרת, בלי docker-compose, בלי שירות שני."),
    ("חלק מהלמידה כאן היא לזהות מתי תשתית היא התשובה הלא נכונה",
     "אם אתם בכל זאת הולכים לשם — לכו בעיניים פקוחות, ותסבירו את השיקול ב-README."),
])

footer_note(s, "עלות",
            "מעבר לקריאות ל-API הפרויקט לא עולה כסף. כל הסטאק חינמי, אין תשתית לשלם עליה.")


# 06 — architecture: classic RAG vs agentic retrieval
s = slide_new(prs)
eyebrow(s, "הארכיטקטורה")
heading(s, "Agentic Retrieval",
        "לא צנרת RAG קלאסית. במקום לשלוף top-k בכל שאלה ולקוות שפגעתם — אתם נותנים ל-Claude כלים, והוא מחליט מה לקרוא.")

cw = (CONTENT_W - 0.35) / 2
bullet_card(s, MARGIN + cw + 0.35, 2.5, cw, 3.85, "הגישה המוכרת", "צנרת RAG קלאסית", [
    "שאלה נכנסת",
    "שולפים top-k קטעים דומים",
    "דוחפים אותם לפרומפט",
    "מקווים שהקטעים הנכונים נכנסו",
], tail="קבוע מראש — רץ אותו דבר בכל שאלה")

bullet_card(s, MARGIN, 2.5, cw, 3.85, "מה שאתם בונים", "כלים על אינדקס מובנה", [
    "שאלה נכנסת",
    "Claude רואה תוכן עניינים וארבעה כלים",
    "הוא מחליט מה לפתוח, ובאיזה סדר",
    "קורא עוד רק אם התשובה עוד לא שלמה",
], dark=True, tail="החלטה בזמן ריצה — לפי מה שהשאלה דורשת")


# 07 — the four tools
s = slide_new(prs)
eyebrow(s, "המימוש")
heading(s, "ארבעת הכלים", "זה מה שאתם חושפים למודל. החיפוש הלקסיקלי עדיין שלכם לממש — פשוט עטוף ככלי.")

tools = [
    ("list_sections()", "תוכן העניינים — מזהה, כותרת, מסמך מקור, עמוד"),
    ("get_section(section_id)", "הטקסט המלא של סעיף אחד"),
    ("search(query, limit=5)", "חיפוש לקסיקלי, BM25 מעל FTS5 — מזהי סעיפים וקטעי הקשר"),
    ("get_course_table(year, semester)", "טבלת קורסים אחת, כשורות מובנות"),
]
y = 2.6
for name, desc in tools:
    rect(s, MARGIN, y, CONTENT_W, 0.78, fill=PAPER, line=HAIRLINE)
    rect(s, MARGIN + CONTENT_W - 0.04, y, 0.04, 0.78, fill=ACCENT)
    text(s, MARGIN + 0.3, y + 0.24, CONTENT_W * 0.52, 0.32, desc, size=13.5, color=INK)
    text(s, MARGIN + CONTENT_W * 0.56, y + 0.24, CONTENT_W * 0.42, 0.32, name,
         size=14, bold=True, font=MONO, color=ACCENT_DARK, align="l", rtl=False)
    y += 0.88

footer_note(s, "עיצוב כלים הוא החלק ההנדסי",
            "כלי שמחזיר יותר מדי מבזבז טוקנים, כלי שמחזיר מעט מדי שולח את המודל לסיבוב נוסף. זה איזון, ומודדים אותו.")


# 08 — why not top-k
s = slide_new(prs)
eyebrow(s, "הנימוק")
heading(s, "למה ככה ולא top-k רגיל", "שלוש סיבות, וכל אחת מהן תפגוש אתכם בשאלות ההערכה.")

reasons = [
    ("01", "שאלות צבירה",
     "״כמה נ״ז סה״כ בשנה ג׳?״ דורשת טבלה שלמה, לא שלושה קטעים דומים סמנטית. "
     "get_course_table מחזיר את כולה."),
    ("02", "המסמכים כבר מובנים",
     "יש סעיפים ממוספרים וטבלאות עם כותרות. לחתוך את זה לחלונות של 1500 תווים "
     "זה להשמיד מידע שקיבלתם בחינם."),
    ("03", "עברית ו-embeddings",
     "זה קרב לא נעים — טוקניזציה, ניקוד, כתיב מלא וחסר. ככה אתם פשוט לא נכנסים אליו."),
]
cw = (CONTENT_W - 0.7) / 3
for i, (num, title, body) in enumerate(reasons):
    cx = MARGIN + (2 - i) * (cw + 0.35)
    rect(s, cx, 2.55, cw, 2.55, fill=PAPER, line=HAIRLINE)
    rect(s, cx + cw - 0.04, 2.55, 0.04, 2.55, fill=ACCENT)
    text(s, cx + 0.28, 2.82, cw - 0.56, 0.26, num, size=11, bold=True, color=ACCENT_DARK,
         font="Georgia", align="l", rtl=False)
    text(s, cx + 0.28, 3.16, cw - 0.56, 0.34, title, size=17, bold=True, color=INK)
    text(s, cx + 0.28, 3.66, cw - 0.56, 1.2, body, size=12, color=MUTED, line_spacing=1.3)

text(s, MARGIN, 5.4, CONTENT_W, 0.5,
     "המבנה של המסמך הוא מידע. צנרת שחותכת אותו לחלונות אחידים זורקת אותו לפח בשלב הראשון.",
     size=14.5, bold=True, color=INK)

footer_note(s, "שימו לב",
            "אתם עדיין מממשים BM25 אמיתי מעל FTS5. ההבדל הוא מי מחליט מתי לקרוא לו — אתם, או המודל.")


# 09 — hard requirements 1+2
s = slide_new(prs)
eyebrow(s, "דרישות חובה  ·  1–2")
heading(s, "מקור וסירוב", "שתי הדרישות שקובעות אם המערכת שלכם ראויה לאמון. פרויקט שלא עומד בהן לא מוגש.")

cw = (CONTENT_W - 0.35) / 2
bullet_card(s, MARGIN + cw + 0.35, 2.4, cw, 3.9, "דרישה 01", "ציטוט מקור בכל תשובה", [
    "מאיזה מסמך נלקחה התשובה",
    "מאיזה סעיף או עמוד",
    "גם כשהתשובה מורכבת משני מקורות",
], tail="תשובה בלי מקור נחשבת שגויה — גם אם היא נכונה")

bullet_card(s, MARGIN, 2.4, cw, 3.9, "דרישה 02", "סירוב כשאין תשובה", [
    "לא ניחוש",
    "לא ידע כללי על מכללות",
    "לא ״כנראה שבערך״",
], dark=True, tail="התשובה היא בדיוק: לא מופיע במסמכים")

footer_note(s, "למה זה נמדד",
            "ההבדל בין עוזר לימודי לבין מקור שקרי הוא בדיוק שתי השורות האלה. בתקנון לימודים, טעות עולה למישהו סמסטר.")


# 10 — hard requirements 3+4
s = slide_new(prs)
eyebrow(s, "דרישות חובה  ·  3–4")
heading(s, "Caching ותקרת איטרציות", "שתי הדרישות ההנדסיות. שתיהן על אותה בעיה: לולאת כלים יקרה.")

rect(s, MARGIN, 2.4, CONTENT_W, 1.9, fill=INK)
text(s, MARGIN + 0.32, 2.62, CONTENT_W - 0.64, 0.22, "דרישה 03  ·  PROMPT CACHING",
     size=10, bold=True, color=ACCENT, spc=300)
text(s, MARGIN + 0.32, 2.92, CONTENT_W - 0.64, 0.34,
     "בלולאת כלים כל פנייה שולחת מחדש את כל מה שקדם לה. בלי cache אתם משלמים על אותו טקסט שוב ושוב.",
     size=14, color=CREAM)
text(s, MARGIN + 0.32, 3.38, CONTENT_W - 0.64, 0.32,
     "cache_control  →  on the last block of every turn",
     size=14, font=MONO, color=ACCENT, align="l", rtl=False)
text(s, MARGIN + 0.32, 3.78, CONTENT_W - 0.64, 0.3,
     "הגדרות הכלים ו-system prompt נשארים תחילית יציבה שנשמרת במטמון בין שאלות.",
     size=12, color=LIGHT_MUTED)

rect(s, MARGIN, 4.5, CONTENT_W, 1.85, fill=PAPER, line=HAIRLINE)
rect(s, MARGIN + CONTENT_W - 0.04, 4.5, 0.04, 1.85, fill=ACCENT)
text(s, MARGIN + 0.32, 4.72, CONTENT_W - 0.64, 0.22, "דרישה 04  ·  תקרת איטרציות",
     size=10, bold=True, color=ACCENT_DARK, spc=300)
text(s, MARGIN + 0.32, 5.02, CONTENT_W - 0.64, 0.34,
     "max_iterations בין 5 ל-8 על לולאת הכלים. סט כלים מעוצב גרוע שולח את המודל חמישה-עשר סיבובים לפני שהוא מוותר.",
     size=14, color=INK)
text(s, MARGIN + 0.32, 5.5, CONTENT_W - 0.64, 0.32,
     "ובנוסף: תדווחו ב-README כמה קריאות כלים בממוצע המערכת עושה לשאלה.",
     size=13.5, bold=True, color=ACCENT_DARK)
text(s, MARGIN + 0.32, 5.9, CONTENT_W - 0.64, 0.3,
     "מדד הנדסי אמיתי, קל למדוד. צוות שהממוצע שלו 9 יודע שיש לו בעיה בעיצוב הכלים.",
     size=12, color=MUTED)

footer_note(s, "החיבור בין השתיים",
            "כל איטרציה מיותרת היא גם טוקנים וגם שנייה של המתנה למשתמש. caching מוזיל אותה, עיצוב כלים טוב מונע אותה.")


# 11 — work plan
s = slide_new(prs)
eyebrow(s, "תכנון")
heading(s, "שלבי העבודה",
        "ארבעה מסלולים שרצים במקביל. שלב 1 הוא החוזה בין כולם — ברגע שיש טקסט מחולץ נקי, השאר עובדים עצמאית.")

rect(s, MARGIN, 2.45, CONTENT_W, 0.85, fill=INK)
text(s, MARGIN + 0.32, 2.62, 3.0, 0.25, "יום 0  ·  ביחד  ·  שעה",
     size=10, bold=True, color=ACCENT, spc=300)
text(s, MARGIN + 0.32, 2.9, CONTENT_W - 0.64, 0.3,
     "תקראו את שני ה-PDF-ים כבני אדם. תחליטו מה סטודנט אמיתי באמת שואל.",
     size=14, color=CREAM)

stages = [
    ("שלב 1", "חילוץ ומבנה", "חילוץ טקסט, תיקון כיווניות, ניקוי תווי בקרה, זיהוי סעיפים וטבלאות, בניית אינדקס FTS5"),
    ("שלב 2", "שכבת המודל", "הגדרת הכלים, לולאת tool use, system prompt, caching, תקרת איטרציות, התנהגות הסירוב"),
    ("שלב 3", "ממשק", "Streamlit, עברית ו-RTL, שאלה נכנסת ← תשובה מבוססת עם ציטוט יוצאת"),
    ("שלב 4", "הערכה", "הרצת עשר שאלות ההערכה, תיעוד הפלטים בפועל, ניתוח הכשלים"),
]
y = 3.55
for tag, title, desc in stages:
    rect(s, MARGIN, y, CONTENT_W, 0.7, fill=PAPER, line=HAIRLINE)
    rect(s, MARGIN + CONTENT_W - 0.04, y, 0.04, 0.7, fill=ACCENT)
    text(s, MARGIN + CONTENT_W - 1.15, y + 0.22, 0.95, 0.28, tag,
         size=11, bold=True, color=ACCENT_DARK)
    text(s, MARGIN + 0.3, y + 0.2, CONTENT_W * 0.62, 0.3, desc, size=12, color=MUTED)
    text(s, MARGIN + CONTENT_W * 0.66, y + 0.18, CONTENT_W * 0.2, 0.32, title,
         size=14.5, bold=True, color=INK)
    y += 0.78

footer_note(s, "שער מעבר בין שלב 1 לשאר",
            "אף מסלול לא מתחיל עד שמישהו פתח את הטקסט המחולץ וקרא אותו. בונים מעל פלט שראיתם, לא פלט שהנחתם.")


# 12 — evaluation questions
s = slide_new(prs)
eyebrow(s, "הערכה")
heading(s, "עשר שאלות, זהות לכל הצוותים", size=36)

table_rtl(s, 1.82,
          ["#", "השאלה", "סוג"],
          [
              ["1", "מהו מספר הקורס של ״מבוא למדעי המחשב״ וכמה נ״ז הוא מזכה?", "עובדתי"],
              ["2", "מה משך הלימודים המקובל לתואר ראשון לפי התקנון?", "עובדתי"],
              ["3", "אילו קורסים נלמדים בשנה ב׳ סמסטר 4, וכמה נ״ז כל אחד?", "טבלה"],
              ["4", "מהן דרישות הקדם של הקורס ״אלגוריתמים 1״?", "טבלה"],
              ["5", "אילו קורסי חובה יש במסלול ״עיבוד אותות ולמידה חישובית״?", "טבלה"],
              ["6", "כמה נ״ז סה״כ נדרשות בשנה ג׳?", "צבירה"],
              ["7", "כמה קורסים בשנה א׳ יש להם דרישות קדם?", "צבירה"],
              ["8", "מה ההבדל בין מסלול חד-חוגי לדו-חוגי, ואיזה רלוונטי למדעי המחשב?", "חוצה מסמכים"],
              ["9", "לפי התקנון, מה התנאים לפריסת לימודים לארבע שנים?", "חוצה מסמכים"],
              ["10", "מהן שעות הקבלה של המזכירות האקדמית?", "לא במסמכים"],
          ],
          widths=[0.6, 9.2, 2.533],
          row_h=0.40,
          sizes=[12, 12.5, 11.5])

footer_note(s, "ההגשה כוללת טבלה עם הפלט בפועל",
            "לא ״עבד״ / ״לא עבד״ — הטקסט שהמערכת שלכם באמת החזירה, לכל אחת מעשר השאלות.")


# 13 — how it is graded
s = slide_new(prs)
eyebrow(s, "מדיניות")
heading(s, "שאלה 10 היא המבחן האמיתי",
        "התשובה הנכונה היחידה לשאלה על שעות קבלה של המזכירות היא שהמידע לא מופיע במסמכים.")

rect(s, MARGIN, 2.5, CONTENT_W, 1.5, fill=INK)
text(s, MARGIN + 0.35, 2.75, CONTENT_W - 0.7, 0.25, "כלל הציון",
     size=10, bold=True, color=ACCENT, spc=300)
text(s, MARGIN + 0.35, 3.1, CONTENT_W - 0.7, 0.5,
     "כישלון גלוי שווה יותר מהצלחה שקטה.", size=26, bold=True, color=CREAM)
text(s, MARGIN + 0.35, 3.62, CONTENT_W - 0.7, 0.3,
     "צוות שמדווח 6 מתוך 10 ומסביר בדיוק למה כל אחת נכשלה — מקבל ציון גבוה יותר מצוות שמדווח 8 בלי ניתוח.",
     size=13, color=LIGHT_MUTED)

row_list(s, 4.25, [
    ("מה נחשב ניתוח כשל", "באיזה שלב זה נשבר: החילוץ? הסעיף שהכלי החזיר? החלטת המודל? ניסוח התשובה?"),
    ("מה לא נחשב ניתוח כשל", "״המודל התבלבל״. זו לא תשובה, זו הודאה שלא הסתכלתם על מה שקרה בלולאה."),
])

footer_note(s, "בשביל זה כדאי ללוגג",
            "שמרו לכל שאלה את רצף קריאות הכלים. בלי הלוג הזה, ניתוח הכשלים יהיה ניחוש — ורואים את זה מיד.")


# 14 — deliverables
s = slide_new(prs)
eyebrow(s, "הגשה")
heading(s, "מה מגישים", "ריפו ציבורי אחד ב-GitHub, ובו שלושה דברים.")

rect(s, MARGIN, 2.4, CONTENT_W, 0.8, fill=PAPER, line=HAIRLINE)
rect(s, MARGIN + CONTENT_W - 0.04, 2.4, 0.04, 0.8, fill=ACCENT)
text(s, MARGIN + CONTENT_W - 1.0, 2.62, 0.8, 0.28, "01", size=11, bold=True,
     color=ACCENT_DARK, font="Georgia")
text(s, MARGIN + 0.3, 2.58, CONTENT_W - 1.4, 0.32, "קוד", size=16, bold=True, color=INK)
text(s, MARGIN + 0.3, 2.9, CONTENT_W - 1.4, 0.28,
     "חילוץ, אינדוקס, כלים, לולאה, ממשק", size=12, color=MUTED)

rect(s, MARGIN, 3.3, CONTENT_W, 2.25, fill=INK)
text(s, MARGIN + CONTENT_W - 1.0, 3.52, 0.8, 0.28, "02", size=11, bold=True,
     color=ACCENT, font="Georgia")
text(s, MARGIN + 0.32, 3.5, CONTENT_W - 1.4, 0.32, "README.md", size=16, bold=True, color=CREAM)
readme_items = [
    "הוראות הרצה שעובדות ממכונה נקייה",
    "פסקה שמנמקת את עיצוב הכלים — למה הסעיפים מפוצלים ככה, מה search מחזיר ולמה",
    "טבלת ההערכה עם הפלטים בפועל וניתוח הכשלים",
    "ממוצע קריאות הכלים לשאלה",
    "השוואה: 3–4 שאלות גם דרך גישה שמכניסה את כל המסמכים ל-context בבת אחת",
]
cy = 3.95
for item in readme_items:
    text(s, MARGIN + 0.72, cy, CONTENT_W - 1.1, 0.28, item, size=12.5, color=LIGHT_MUTED)
    text(s, MARGIN + CONTENT_W - 0.64, cy + 0.01, 0.18, 0.22, "·", size=16, bold=True, color=ACCENT)
    cy += 0.32

rect(s, MARGIN, 5.65, CONTENT_W, 0.8, fill=PAPER, line=HAIRLINE)
rect(s, MARGIN + CONTENT_W - 0.04, 5.65, 0.04, 0.8, fill=ACCENT)
text(s, MARGIN + CONTENT_W - 1.0, 5.87, 0.8, 0.28, "03", size=11, bold=True,
     color=ACCENT_DARK, font="Georgia")
text(s, MARGIN + 0.3, 5.83, CONTENT_W - 1.4, 0.32, "הדגמה", size=16, bold=True, color=INK)
text(s, MARGIN + 0.3, 6.15, CONTENT_W - 1.4, 0.28,
     "סרטון של 3 דקות, או הדגמה חיה בכיתה", size=12, color=MUTED)

footer_note(s, "ההשוואה בסעיף 2 עולה שעה והיא המשפט הכי בעל ערך בהגשה",
            "כל הקורפוס נכנס לחלון הקשר אחד. כדאי שתדעו מה בדיוק קניתם בזה שבניתם שליפה.", y=6.6)


# 15 — API keys and budget
s = slide_new(prs, dark=True)
text(s, MARGIN, 0.35, CONTENT_W, 0.3, "כסף  ·  קראו את זה עד הסוף", size=11, bold=True,
     color=ACCENT, spc=400)
heading(s, "מפתחות API ותקציב",
        "מפתח ה-API של הפרויקט מגיע ממני, עם תקרת הוצאה לכל צוות.",
        color=CREAM, sub_color=LIGHT_MUTED)

warns = [
    ("01", "המפתח הזה הוא לאפליקציה שלכם — לא ל-Claude Code",
     "אל תגדירו אותו כ-ANTHROPIC_API_KEY גלובלי. Claude Code שורף סדר גודל יותר טוקנים "
     "מהצ׳אטבוט שתבנו, והוא ירוקן את התקרה תוך שעות. ל-Claude Code תשתמשו במנוי שלכם."),
    ("02", "לא מעלים לאוויר",
     "אפליקציית Streamlit ציבורית עם המפתח שלי ובלי אימות היא דרך לאבד כסף אמיתי בלילה אחד. "
     "רץ מקומית: streamlit run, המפתח ב-.env, ו-.env ב-.gitignore."),
    ("03", "תעקבו אחרי מה שאתם שורפים",
     "התקרה מגינה עליי, לא עליכם. צוות שמגיע אליה ביום השני נשאר בלי מפתח."),
]
y = 2.5
for num, title, body in warns:
    rect(s, MARGIN, y, CONTENT_W, 1.3, fill=INK_SOFT)
    rect(s, MARGIN + CONTENT_W - 0.04, y, 0.04, 1.3, fill=ACCENT)
    text(s, MARGIN + CONTENT_W - 1.0, y + 0.22, 0.8, 0.26, num, size=11, bold=True,
         color=ACCENT, font="Georgia")
    text(s, MARGIN + 0.32, y + 0.2, CONTENT_W - 1.4, 0.32, title, size=15.5, bold=True, color=CREAM)
    text(s, MARGIN + 0.32, y + 0.62, CONTENT_W - 0.64, 0.55, body, size=12, color=LIGHT_MUTED,
         line_spacing=1.25)
    y += 1.42

text(s, MARGIN, 6.85, CONTENT_W, 0.3,
     "מעבר לזה הפרויקט לא עולה כסף. כל הסטאק חינמי, אין תשתית לשלם עליה, ואין ספק שני.",
     size=12, color=DIM)


# 16 — Hebrew is the project
s = slide_new(prs)
eyebrow(s, "הערה שתחסוך לכם ערב")
heading(s, "עברית היא לא דקורציה — היא הפרויקט",
        "היא מקור אמיתי של קושי הנדסי, בשלושה מקומות שונים לפחות. אתם תיתקלו בהם. זו לא תקלה בפרויקט.")

places = [
    ("בחילוץ", "סדר התווים שיוצא מ-pdfplumber לא בהכרח סדר הקריאה"),
    ("באינדוקס", "טוקניזציה של עברית — כתיב מלא וחסר, מקף, גרשיים"),
    ("בתצוגה", "RTL בדפדפן, ומחרוזות מעורבות עברית-אנגלית-מספרים"),
]
cw = (CONTENT_W - 0.7) / 3
for i, (title, body) in enumerate(places):
    cx = MARGIN + (2 - i) * (cw + 0.35)
    rect(s, cx, 2.5, cw, 1.5, fill=PAPER, line=HAIRLINE)
    rect(s, cx + cw - 0.04, 2.5, 0.04, 1.5, fill=ACCENT)
    text(s, cx + 0.28, 2.76, cw - 0.56, 0.34, title, size=17, bold=True, color=INK)
    text(s, cx + 0.28, 3.22, cw - 0.56, 0.7, body, size=12, color=MUTED, line_spacing=1.3)

rect(s, MARGIN, 4.3, CONTENT_W, 1.85, fill=INK)
text(s, MARGIN + 0.35, 4.55, CONTENT_W - 0.7, 0.25, "כשמשהו נראה כמו ג׳יבריש",
     size=10, bold=True, color=ACCENT, spc=300)
text(s, MARGIN + 0.35, 4.88, CONTENT_W - 0.7, 0.35,
     "תעצרו, תסתכלו על הבייטים, ותבינו מה קרה — לפני שאתם כותבים תיקון.",
     size=17, bold=True, color=CREAM)
text(s, MARGIN + 0.35, 5.4, CONTENT_W - 0.7, 0.32,
     "s[::-1]  # ← זה לא תיקון. זה באג חדש.",
     size=14, font=MONO, color=ACCENT, align="l", rtl=False)
text(s, MARGIN + 0.35, 5.78, CONTENT_W - 0.7, 0.3,
     "היפוך מחרוזת מעורבת עברית-אנגלית-מספרים שובר בדיוק את מה שהיה תקין.",
     size=12, color=LIGHT_MUTED)

footer_note(s, "נתקעתם?",
            "תביאו את השגיאה למפגש. חצי מהקורס הוא לראות איך מפרקים תקלה אמיתית — וכאן הן יהיו אמיתיות.")


# 17 — closing
s = slide_new(prs, dark=True)
rect(s, MARGIN + CONTENT_W - 1.2, 3.66, 1.2, 0.04, fill=ACCENT)

text(s, MARGIN, 1.5, CONTENT_W, 0.4, "מתחילים", size=13, bold=True, color=ACCENT, spc=400)
text(s, MARGIN, 1.92, CONTENT_W, 0.9, "בהצלחה", size=58, bold=True, color=CREAM)
text(s, MARGIN, 3.0, CONTENT_W, 0.5,
     "הבריף המלא, שני ה-PDF-ים וסט השאלות — הכול בריפו של הקורס, בתיקייה final-project.",
     size=17, color=LIGHT_MUTED)

links = [
    ("הבריף המלא של הפרויקט", "claude101/final-project/README.md",
     "https://github.com/talgurevich/claude101/tree/main/final-project"),
    ("הריפו של הקורס", "github.com/talgurevich/claude101",
     "https://github.com/talgurevich/claude101"),
    ("מייל", "tal.gurevich@gmail.com", "mailto:tal.gurevich@gmail.com"),
]
y = 4.1
for label, shown, href in links:
    text(s, MARGIN + CONTENT_W * 0.5, y, CONTENT_W * 0.5, 0.3, label, size=13, color=CREAM)
    text(s, MARGIN, y, CONTENT_W * 0.48, 0.3,
         [(shown, {"link": href, "color": ACCENT})],
         size=13, font=MONO, color=ACCENT, align="l", rtl=False)
    rect(s, MARGIN, y + 0.4, CONTENT_W, 0.014, fill=HAIRLINE_DARK)
    y += 0.62

text(s, MARGIN, 6.5, CONTENT_W, 0.4, "המכללה האקדמית תל-חי  ·  טל גורביץ׳", size=12, color=DIM)


out = "/Users/talgurevich/Documents/Claude101/final-project/final-project-deck-he.pptx"
prs.save(out)
print("wrote", out, "|", len(prs.slides._sldIdLst), "slides")
